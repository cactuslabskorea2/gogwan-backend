from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import Response
from pydantic import BaseModel
from starlette.middleware.base import BaseHTTPMiddleware
import base64
import os
import glob
from google import genai
from google.genai import types
from google.genai.types import GenerateContentConfig, Part, Modality
from korean_lunar_calendar import KoreanLunarCalendar
from datetime import datetime
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
import json
from google.cloud import firestore

app = FastAPI(title="GOGWAN API")

# CSP 미들웨어 추가
class CSPMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        # HTML 파일에만 CSP 헤더 추가
        if response.headers.get("content-type", "").startswith("text/html"):
            csp_policy = (
                "default-src 'self'; "
                "script-src 'self' 'unsafe-inline' 'unsafe-eval' https://www.gstatic.com https://*.firebaseio.com https://cdn.jsdelivr.net https://generativelanguage.googleapis.com https://*.googleapis.com; "
                "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
                "img-src 'self' data: blob: https://*.googleusercontent.com https://firebasestorage.googleapis.com https://*.firebaseapp.com; "
                "font-src 'self' data: https://fonts.gstatic.com; "
                "connect-src 'self' "
                "https://generativelanguage.googleapis.com "
                "https://firestore.googleapis.com "
                "https://*.googleapis.com "
                "https://*.firebaseio.com "
                "https://www.gstatic.com "
                "https://*.run.app;"
            )
            response.headers["Content-Security-Policy"] = csp_policy
        return response

app.add_middleware(CSPMiddleware)

# 스타일 설명 캐시 (메모리에 저장)
style_descriptions = {}

# CORS 설정
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Vertex AI 클라이언트 (나노바나나 사용)
# Cloud Run 환경 감지
is_cloud_run = os.getenv('K_SERVICE') is not None

if is_cloud_run:
    # Cloud Run 환경: 기본 인증 사용 (Service Account 자동 적용)
    pass
elif os.getenv('GOOGLE_APPLICATION_CREDENTIALS_JSON'):
    import json
    import tempfile
    # 환경 변수에서 JSON 읽어서 임시 파일 생성 (Railway용)
    creds_json = json.loads(os.getenv('GOOGLE_APPLICATION_CREDENTIALS_JSON'))
    with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.json') as f:
        json.dump(creds_json, f)
        os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = f.name
else:
    # 로컬 환경: 파일 사용
    credentials_path = os.path.join(
        os.path.dirname(__file__),
        'gogwan-4902b9a702be.json'
    )
    if os.path.exists(credentials_path):
        os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = credentials_path

client = genai.Client(
    vertexai=True,
    project='gogwan',
    location='global'
)

# Firestore 클라이언트 초기화 (lazy loading)
db = None

def get_firestore_client():
    global db
    if db is None:
        import os
        cred_path = os.path.join(os.path.dirname(__file__), 'gogwan-4902b9a702be.json')
        os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = cred_path
        # Use gogwan-e79bc project (matches google-services.json in frontend)
        # Even though service account is from 'gogwan' project, it should work cross-project
        db = firestore.Client(project='gogwan-e79bc')
    return db


async def analyze_reference_image(image_path: str, style_id: str):
    """레퍼런스 이미지를 분석해서 스타일 설명을 생성"""
    try:
        with open(image_path, 'rb') as f:
            image_bytes = f.read()

        mime_type = "image/png" if image_path.endswith('.png') else "image/jpeg"

        prompt = """Analyze this image in detail and describe its artistic style for the purpose of recreating similar images.

Describe:
1. Art style (e.g., anime, cartoon, realistic painting, watercolor, oil painting, digital art, etc.)
2. Color palette and color grading (warm/cool tones, saturation, contrast)
3. Lighting style and mood (soft, dramatic, natural, studio lighting)
4. Texture and rendering technique (smooth, painterly, sketchy, 3D rendered)
5. Background style and atmosphere
6. Overall aesthetic and visual characteristics

Provide a concise but comprehensive description that can be used to recreate images in this exact style."""

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                Part.from_bytes(data=image_bytes, mime_type=mime_type),
                prompt
            ]
        )

        description = response.text
        style_descriptions[style_id] = description
        print(f"✅ Analyzed style '{style_id}': {description[:100]}...")
        return description

    except Exception as e:
        print(f"❌ Failed to analyze style '{style_id}': {str(e)}")
        return None


@app.on_event("startup")
async def startup_event():
    """서버 시작 시 모든 레퍼런스 이미지 분석"""
    print("🚀 Starting up... Analyzing reference images...")

    reference_dir = os.path.join(os.path.dirname(__file__), 'reference_images')

    image_files = []
    for ext in ['*.jpg', '*.jpeg', '*.png']:
        image_files.extend(glob.glob(os.path.join(reference_dir, ext)))

    for image_path in image_files:
        filename = os.path.basename(image_path)
        style_id = os.path.splitext(filename)[0]

        if style_id.startswith('.') or style_id.lower() == 'readme':
            continue

        await analyze_reference_image(image_path, style_id)

    print(f"✅ Analyzed {len(style_descriptions)} styles")


class IdPhotoRequest(BaseModel):
    image: str  # 사진 (필수)
    background_color: str = "white"


class StyleTransferRequest(BaseModel):
    image: str
    style: str = None  # cartoon, western, medieval, hanbok 등 (optional)
    reference_image: str = None  # 커스텀 레퍼런스 이미지 (Base64, optional)


class PressReleaseRequest(BaseModel):
    prompt: str
    reference_content: str = None


@app.get("/")
def read_root():
    return {"message": "GOGWAN API Server", "status": "running"}


@app.get("/api/available-styles")
async def get_available_styles():
    """사용 가능한 스타일 목록과 썸네일 반환"""
    try:
        reference_dir = os.path.join(os.path.dirname(__file__), 'reference_images')

        # jpg, jpeg, png 파일 찾기
        image_files = []
        for ext in ['*.jpg', '*.jpeg', '*.png']:
            image_files.extend(glob.glob(os.path.join(reference_dir, ext)))

        styles = []
        for image_path in image_files:
            # 파일명에서 스타일 ID 추출 (확장자 제외)
            filename = os.path.basename(image_path)
            style_id = os.path.splitext(filename)[0]

            # README.md, .gitkeep 등 제외
            if style_id.startswith('.') or style_id.lower() == 'readme':
                continue

            # 이미지를 base64로 인코딩
            with open(image_path, 'rb') as f:
                image_bytes = f.read()
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')

            # MIME 타입 결정
            mime_type = 'image/png' if image_path.endswith('.png') else 'image/jpeg'

            styles.append({
                "id": style_id,
                "thumbnail": f"data:{mime_type};base64,{image_base64}"
            })

        return {
            "success": True,
            "styles": styles
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"스타일 목록을 불러오는 중 오류 발생: {str(e)}")


@app.post("/api/create-id-photo")
async def create_id_photo(request: IdPhotoRequest):
    """증명사진 생성 - Gemini 2.5 Flash Image"""
    try:
        # 이미지 Base64 디코딩
        image_bytes = base64.b64decode(request.image)

        # 프롬프트
        prompt = f"""Generate a professional ID photo from this photograph.

Requirements:
1. ASPECT RATIO: 7:9 (width:height) - Standard Korean ID photo ratio (3.5cm × 4.5cm)
2. The person facing directly forward (front-facing), looking straight at the camera
3. Dressed in formal business attire (suit and tie for men, formal blouse or suit for women)
4. Sitting upright with proper posture, shoulders straight
5. Natural, professional facial expression with a slight, confident smile
6. Studio-quality lighting: soft, even, professional lighting on the face
7. Clean, solid {request.background_color} background
8. Standard ID photo composition: head and upper shoulders visible, face centered
9. Head length (crown to chin) should be approximately 70-80% of the image height
10. High resolution and clarity suitable for passports, national ID cards, and resumes
11. Professional photo studio quality

IMPORTANT: The output image MUST be in 7:9 aspect ratio (portrait orientation).
Preserve the person's facial features, proportions, and identity exactly."""

        # Gemini 2.5 Flash Image (나노바나나)로 이미지 생성
        response = client.models.generate_content(
            model="gemini-2.5-flash-image",
            contents=[
                prompt,
                Part.from_bytes(data=image_bytes, mime_type="image/jpeg")
            ],
            config=GenerateContentConfig(
                response_modalities=[Modality.IMAGE]
            )
        )

        # 생성된 이미지 추출
        for part in response.candidates[0].content.parts:
            if part.inline_data is not None:
                result_image_base64 = base64.b64encode(part.inline_data.data).decode('utf-8')
                return {
                    "success": True,
                    "message": "증명사진이 생성되었습니다",
                    "processed_image": result_image_base64
                }

        raise HTTPException(status_code=500, detail="이미지 생성에 실패했습니다")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"처리 중 오류 발생: {str(e)}")


@app.post("/api/convert-to-style")
async def convert_to_style(request: StyleTransferRequest):
    """스타일 변환 - Gemini 2.5 Flash Image with Multiple Images"""
    try:
        # 사용자 이미지 디코딩
        user_image_bytes = base64.b64decode(request.image)

        # 레퍼런스 이미지 처리
        if request.reference_image:
            # 커스텀 레퍼런스 이미지 사용
            reference_image_bytes = base64.b64decode(request.reference_image)
            reference_mime_type = "image/jpeg"  # 기본값
        elif request.style:
            # 저장된 스타일 사용
            reference_image_path = os.path.join(
                os.path.dirname(__file__),
                'reference_images',
                f'{request.style}.jpg'
            )

            if not os.path.exists(reference_image_path):
                reference_image_path = reference_image_path.replace('.jpg', '.png')

            if not os.path.exists(reference_image_path):
                raise HTTPException(
                    status_code=404,
                    detail=f"'{request.style}' 레퍼런스 이미지를 찾을 수 없습니다."
                )

            with open(reference_image_path, 'rb') as f:
                reference_image_bytes = f.read()

            reference_mime_type = "image/png" if reference_image_path.endswith('.png') else "image/jpeg"
        else:
            raise HTTPException(
                status_code=400,
                detail="스타일 ID 또는 커스텀 레퍼런스 이미지를 제공해주세요."
            )

        # 프롬프트 - Google 공식 문서 베스트 프랙티스 기반
        prompt = """Transform the provided photograph into the artistic style shown in the reference image.

PRESERVE THE ORIGINAL COMPOSITION: Keep the person's exact pose, body position, facial expression, and overall layout exactly as shown in the photograph.

APPLY THE ARTISTIC STYLE: Render the entire scene using the artistic techniques, visual characteristics, and aesthetic from the reference image, including:
- Art style and rendering technique (illustration style, painting technique, digital art approach)
- Color palette and color grading
- Lighting mood and atmosphere
- Texture and brushwork style
- Line art and shading approach
- Overall visual aesthetic

The person and their pose should remain identical to the original photograph, but completely redrawn/repainted using the artistic style from the reference image. The background should also be created in the same artistic style to create a cohesive, unified artwork."""

        response = client.models.generate_content(
            model="gemini-2.5-flash-image",
            contents=[
                Part.from_bytes(data=reference_image_bytes, mime_type=reference_mime_type),
                Part.from_bytes(data=user_image_bytes, mime_type="image/jpeg"),
                prompt
            ],
            config=GenerateContentConfig(
                response_modalities=[Modality.IMAGE]
            )
        )

        for part in response.candidates[0].content.parts:
            if part.inline_data is not None:
                result_image_base64 = base64.b64encode(part.inline_data.data).decode('utf-8')

                # 메시지 생성
                if request.reference_image:
                    message = "커스텀 스타일 이미지가 생성되었습니다"
                else:
                    message = f"{request.style} 스타일 이미지가 생성되었습니다"

                return {
                    "success": True,
                    "message": message,
                    "processed_image": result_image_base64
                }

        raise HTTPException(status_code=500, detail="이미지 생성에 실패했습니다")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"처리 중 오류 발생: {str(e)}")


class BackgroundRemovalRequest(BaseModel):
    image: str  # Base64 인코딩된 이미지


@app.post("/api/remove-background")
async def remove_background(request: BackgroundRemovalRequest):
    """배경 제거 - Gemini 2.5 Flash Image"""
    try:
        # 이미지 Base64 디코딩
        image_bytes = base64.b64decode(request.image)

        # 프롬프트
        prompt = """Remove the background from this image completely.

Requirements:
1. Keep the main subject (person, object, etc.) exactly as it is
2. Remove ALL background elements completely
3. Make the background fully transparent (PNG with alpha channel)
4. Preserve all details of the main subject
5. Clean, precise edges around the subject
6. Professional quality background removal
7. Do not modify or change the main subject in any way

Output must be a PNG image with transparent background."""

        # Gemini 2.5 Flash Image로 배경 제거
        response = client.models.generate_content(
            model="gemini-2.5-flash-image",
            contents=[
                prompt,
                Part.from_bytes(data=image_bytes, mime_type="image/jpeg")
            ],
            config=GenerateContentConfig(
                response_modalities=[Modality.IMAGE]
            )
        )

        # 생성된 이미지 추출
        for part in response.candidates[0].content.parts:
            if part.inline_data is not None:
                result_image_base64 = base64.b64encode(part.inline_data.data).decode('utf-8')
                return {
                    "success": True,
                    "message": "배경이 제거되었습니다",
                    "processed_image": result_image_base64
                }

        raise HTTPException(status_code=500, detail="배경 제거에 실패했습니다")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"처리 중 오류 발생: {str(e)}")


# 하위 호환성을 위한 기존 엔드포인트 유지
@app.post("/api/convert-to-ghibli")
async def convert_to_ghibli_legacy(request: StyleTransferRequest):
    """지브리풍 변환 (레거시 엔드포인트 - convert-to-style 사용 권장)"""
    request.style = "cartoon"
    return await convert_to_style(request)


@app.post("/api/generate-press-release")
async def generate_press_release(request: PressReleaseRequest):
    """보도자료 생성"""
    try:
        full_prompt = '당신은 대한민국 정부기관의 전문 보도자료 작성자입니다.\n공식적이고 명확하며, 객관적인 톤으로 작성해주세요.\n\n'

        if request.reference_content:
            full_prompt += f'다음은 참고할 기존 보도자료입니다:\n\n{request.reference_content}\n\n위 보도자료의 스타일과 구조를 참고하여 작성해주세요.\n\n'

        full_prompt += f'''사용자 요청사항:
{request.prompt}

다음 구조로 대한민국 정부기관 공식 보도자료 형식에 맞춰 작성해주세요:

【제목】
간결하고 핵심을 담은 한 줄 제목

【부제】
제목을 보완하는 설명 (2줄 이내)

【본문】
본문은 자연스러운 문장으로 작성하되, 다음 구조를 따르세요:

□ 도입부
  ○ 정책/사업의 배경과 목적을 간단히 설명

□ 핵심 내용
  ○ 지원 대상, 금액, 기간, 방법 등 주요 정보 포함
  ○ 숫자나 통계는 구체적으로 표기

□ 향후 계획 및 기대효과
  ○ 기대효과와 향후 계획을 명시

작성 규칙:
- □와 ○ 기호만 사용하고, 그 뒤에 바로 내용을 작성하세요
- 각 문단은 자연스러운 문장으로 구성
- 공식적이고 전문적인 톤을 유지

보도자료는 반드시 한국어로 작성하며, 【제목】【부제】【본문】 구조를 명확히 지켜주세요.'''

        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=full_prompt
        )

        text = response.text

        if not text:
            raise HTTPException(status_code=500, detail="Gemini API가 빈 응답을 반환했습니다")

        return {"success": True, "content": text}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"처리 중 오류 발생: {str(e)}")


# ============================================
# 사주팔자 계산 로직
# ============================================

# 천간 (Heavenly Stems)
HEAVENLY_STEMS = ["갑", "을", "병", "정", "무", "기", "경", "신", "임", "계"]
# 지지 (Earthly Branches)
EARTHLY_BRANCHES = ["자", "축", "인", "묘", "진", "사", "오", "미", "신", "유", "술", "해"]

# 오행 (Five Elements)
FIVE_ELEMENTS = {
    "갑": "목", "을": "목",
    "병": "화", "정": "화",
    "무": "토", "기": "토",
    "경": "금", "신": "금",
    "임": "수", "계": "수",
    "인": "목", "묘": "목",
    "사": "화", "오": "화",
    "진": "토", "술": "토", "축": "토", "미": "토",
    "신": "금", "유": "금",
    "자": "수", "해": "수"
}

# 음양
YIN_YANG = {
    "갑": "양", "을": "음",
    "병": "양", "정": "음",
    "무": "양", "기": "음",
    "경": "양", "신": "음",
    "임": "양", "계": "음",
    "자": "양", "축": "음",
    "인": "양", "묘": "음",
    "진": "양", "사": "음",
    "오": "양", "미": "음",
    "신": "양", "유": "음",
    "술": "양", "해": "음"
}

# 십신 (Ten Gods)
def get_sipsin(day_stem, target_stem):
    """일간과 대상 간지로 십신 계산"""
    day_idx = HEAVENLY_STEMS.index(day_stem)
    target_idx = HEAVENLY_STEMS.index(target_stem)

    day_element = FIVE_ELEMENTS[day_stem]
    target_element = FIVE_ELEMENTS[target_stem]
    day_yang = YIN_YANG[day_stem]
    target_yang = YIN_YANG[target_stem]

    # 나를 생하는 것 (생아): 정인, 편인
    # 내가 생하는 것 (아생): 식신, 상관
    # 나를 극하는 것 (극아): 정관, 편관(칠살)
    # 내가 극하는 것 (아극): 정재, 편재
    # 나와 같은 것 (비화): 비견, 겁재

    element_relation = get_element_relation(day_element, target_element)
    same_yang = (day_yang == target_yang)

    if element_relation == "same":
        return "비견" if same_yang else "겁재"
    elif element_relation == "generate":
        return "식신" if same_yang else "상관"
    elif element_relation == "control":
        return "편재" if same_yang else "정재"
    elif element_relation == "controlled":
        return "편관" if same_yang else "정관"
    elif element_relation == "generated":
        return "편인" if same_yang else "정인"

    return "비견"

def get_element_relation(from_elem, to_elem):
    """오행 상생상극 관계"""
    generate_cycle = {"목": "화", "화": "토", "토": "금", "금": "수", "수": "목"}
    control_cycle = {"목": "토", "토": "수", "수": "화", "화": "금", "금": "목"}

    if from_elem == to_elem:
        return "same"
    elif generate_cycle.get(from_elem) == to_elem:
        return "generate"
    elif control_cycle.get(from_elem) == to_elem:
        return "control"
    elif generate_cycle.get(to_elem) == from_elem:
        return "generated"
    elif control_cycle.get(to_elem) == from_elem:
        return "controlled"

    return "same"


def calculate_saju(year, month, day, hour, is_lunar=False):
    """사주팔자 계산"""

    # 양력을 음력으로 변환 (필요시)
    if not is_lunar:
        calendar = KoreanLunarCalendar()
        calendar.setSolarDate(year, month, day)
        lunar_year = calendar.lunarYear
        lunar_month = calendar.lunarMonth
        lunar_day = calendar.lunarDay
        is_intercalation = calendar.isIntercalation
    else:
        lunar_year = year
        lunar_month = month
        lunar_day = day
        is_intercalation = False

    # 년주 계산 (입춘 기준)
    year_stem_idx = (lunar_year - 4) % 10
    year_branch_idx = (lunar_year - 4) % 12
    year_pillar = HEAVENLY_STEMS[year_stem_idx] + EARTHLY_BRANCHES[year_branch_idx]

    # 월주 계산
    month_stem_idx = (year_stem_idx * 2 + lunar_month + 1) % 10
    month_branch_idx = (lunar_month + 1) % 12
    month_pillar = HEAVENLY_STEMS[month_stem_idx] + EARTHLY_BRANCHES[month_branch_idx]

    # 일주 계산 (만세력 기준일로부터 계산)
    # 기준일: 1900년 1월 1일 = 갑자일
    base_date = datetime(1900, 1, 1)
    if not is_lunar:
        target_date = datetime(year, month, day)
    else:
        calendar = KoreanLunarCalendar()
        calendar.setLunarDate(lunar_year, lunar_month, lunar_day, is_intercalation)
        target_date = datetime(calendar.solarYear, calendar.solarMonth, calendar.solarDay)

    days_diff = (target_date - base_date).days
    day_stem_idx = (days_diff + 10) % 10
    day_branch_idx = (days_diff + 0) % 12
    day_pillar = HEAVENLY_STEMS[day_stem_idx] + EARTHLY_BRANCHES[day_branch_idx]

    # 시주 계산
    hour_branch_idx = ((hour + 1) // 2) % 12
    hour_stem_idx = (day_stem_idx * 2 + hour_branch_idx) % 10
    hour_pillar = HEAVENLY_STEMS[hour_stem_idx] + EARTHLY_BRANCHES[hour_branch_idx]

    return {
        "year_pillar": year_pillar,
        "month_pillar": month_pillar,
        "day_pillar": day_pillar,
        "hour_pillar": hour_pillar,
        "day_stem": HEAVENLY_STEMS[day_stem_idx]
    }


def analyze_ohaeng(saju):
    """오행 분석"""
    pillars = [saju["year_pillar"], saju["month_pillar"], saju["day_pillar"], saju["hour_pillar"]]
    elements = {"목": 0, "화": 0, "토": 0, "금": 0, "수": 0}

    for pillar in pillars:
        for char in pillar:
            if char in FIVE_ELEMENTS:
                elements[FIVE_ELEMENTS[char]] += 1

    return elements


def calculate_daeun(birth_year, birth_month, birth_day, year_pillar, gender):
    """대운 계산"""
    # year_pillar에서 천간과 지지 추출
    year_stem = year_pillar[0]
    year_branch = year_pillar[1]

    # 양남음녀: 순행, 음남양녀: 역행
    is_yang_year = YIN_YANG[year_stem] == "양"
    is_male = (gender == "남")

    forward = (is_yang_year and is_male) or (not is_yang_year and not is_male)

    # 대운 시작 나이 계산 (간략화)
    start_age = 8  # 평균적으로 8세 전후

    daeun_list = []
    year_stem_idx = HEAVENLY_STEMS.index(year_stem)
    year_branch_idx = EARTHLY_BRANCHES.index(year_branch)

    for i in range(8):  # 8개 대운
        age = start_age + (i * 10)
        if forward:
            stem_idx = (year_stem_idx + i + 1) % 10
            branch_idx = (year_branch_idx + i + 1) % 12
        else:
            stem_idx = (year_stem_idx - i - 1) % 10
            branch_idx = (year_branch_idx - i - 1) % 12

        daeun = HEAVENLY_STEMS[stem_idx] + EARTHLY_BRANCHES[branch_idx]
        daeun_list.append({"age": age, "pillar": daeun})

    return daeun_list


class SajuLifetimeRequest(BaseModel):
    birth_year: int
    birth_month: int
    birth_day: int
    birth_hour: int
    gender: str  # "남" 또는 "여"
    is_lunar: bool = False


class SajuYearlyRequest(BaseModel):
    birth_year: int
    birth_month: int
    birth_day: int
    birth_hour: int
    gender: str
    target_year: int  # 운세를 볼 년도
    is_lunar: bool = False


@app.post("/api/saju-lifetime")
async def saju_lifetime(request: SajuLifetimeRequest):
    """평생운 사주 해석"""
    try:
        # 사주 계산
        saju = calculate_saju(
            request.birth_year,
            request.birth_month,
            request.birth_day,
            request.birth_hour,
            request.is_lunar
        )

        # 오행 분석
        ohaeng = analyze_ohaeng(saju)

        # 대운 계산
        daeun = calculate_daeun(
            request.birth_year,
            request.birth_month,
            request.birth_day,
            saju["year_pillar"],
            request.gender
        )

        # Gemini API로 해석 생성
        prompt = f"""당신은 전문 명리학자입니다. 다음 사주를 바탕으로 평생운을 자세히 풀이해주세요.

【사주팔자】
년주: {saju["year_pillar"]}
월주: {saju["month_pillar"]}
일주: {saju["day_pillar"]} (일간: {saju["day_stem"]})
시주: {saju["hour_pillar"]}

【오행 분석】
목: {ohaeng["목"]}개, 화: {ohaeng["화"]}개, 토: {ohaeng["토"]}개, 금: {ohaeng["금"]}개, 수: {ohaeng["수"]}개

【대운】
{", ".join([f"{d['age']}세({d['pillar']})" for d in daeun[:4]])}

다음 항목으로 풀이해주세요:

1. 타고난 성격과 기질
2. 재물운과 사업운
3. 관운(직장운)과 명예운
4. 건강운
5. 인간관계와 대인운
6. 결혼운과 가정운
7. 말년운

각 항목마다 구체적이고 상세하게 설명해주세요."""

        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=prompt
        )

        interpretation = response.text

        return {
            "success": True,
            "saju": saju,
            "ohaeng": ohaeng,
            "daeun": daeun,
            "interpretation": interpretation
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"사주 해석 중 오류 발생: {str(e)}")


@app.post("/api/saju-yearly")
async def saju_yearly(request: SajuYearlyRequest):
    """신년운 사주 해석"""
    try:
        # 사주 계산
        saju = calculate_saju(
            request.birth_year,
            request.birth_month,
            request.birth_day,
            request.birth_hour,
            request.is_lunar
        )

        # 해당 년도의 세운 계산
        year_stem_idx = (request.target_year - 4) % 10
        year_branch_idx = (request.target_year - 4) % 12
        saeun = HEAVENLY_STEMS[year_stem_idx] + EARTHLY_BRANCHES[year_branch_idx]

        # 오행 분석
        ohaeng = analyze_ohaeng(saju)

        # Gemini API로 해석 생성
        prompt = f"""당신은 전문 명리학자입니다. {request.target_year}년 신년운을 자세히 풀이해주세요.

【사주팔자】
년주: {saju["year_pillar"]}
월주: {saju["month_pillar"]}
일주: {saju["day_pillar"]} (일간: {saju["day_stem"]})
시주: {saju["hour_pillar"]}

【{request.target_year}년 세운】
{saeun}년

【오행 분석】
목: {ohaeng["목"]}개, 화: {ohaeng["화"]}개, 토: {ohaeng["토"]}개, 금: {ohaeng["금"]}개, 수: {ohaeng["수"]}개

다음 항목으로 {request.target_year}년 운세를 풀이해주세요:

1. 총운 (전반적인 운세)
2. 재물운 (재테크, 투자, 수입)
3. 사업운 (직장, 승진, 이직)
4. 건강운 (주의할 질병, 건강 관리)
5. 애정운 (연애, 결혼, 가족관계)
6. 학업/시험운
7. 월별 운세 흐름 (상반기/하반기)

각 항목마다 구체적이고 실용적인 조언을 포함해서 설명해주세요."""

        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=prompt
        )

        interpretation = response.text

        return {
            "success": True,
            "saju": saju,
            "saeun": saeun,
            "target_year": request.target_year,
            "ohaeng": ohaeng,
            "interpretation": interpretation
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"신년운 해석 중 오류 발생: {str(e)}")


# ============================================
# 현수막 생성
# ============================================

class BannerRequest(BaseModel):
    message: str  # 현수막 주 문구
    date: str = ""  # 날짜/시간
    location: str = ""  # 장소
    organizer: str = ""  # 주최
    host: str = ""  # 주관
    width_meter: float = 5.0  # 가로 크기 (m)
    height_meter: float = 1.0  # 세로 크기 (m)


@app.post("/api/generate-banner")
async def generate_banner(request: BannerRequest):
    """
    현수막 배경 생성 (Gemini) + 텍스트 오버레이 (Pillow)
    - 다운로드: 깨끗한 배경 이미지만 제공
    - 미리보기: 배경 + 텍스트 합성
    """
    try:
        # 미터를 픽셀로 변환 (인쇄용 고해상도: 300 DPI 기준)
        # 1 미터 = 39.37 인치, 300 DPI 기준
        dpi = 150  # 웹 미리보기용 적절한 해상도
        pixels_per_meter = int(39.37 * dpi)

        width_px = int(request.width_meter * pixels_per_meter)
        height_px = int(request.height_meter * pixels_per_meter)

        # 최대 크기 제한 (메모리 절약)
        max_dimension = 6000
        if width_px > max_dimension or height_px > max_dimension:
            ratio = min(max_dimension / width_px, max_dimension / height_px)
            width_px = int(width_px * ratio)
            height_px = int(height_px * ratio)

        # 1단계: Gemini로 주제에 맞는 배경 디자인 생성
        prompt = f"""Create a professional Korean banner (현수막) background design WITHOUT ANY TEXT.

Banner topic/theme: "{request.message}"

Design specifications:
- Size: {width_px}x{height_px}px (aspect ratio {request.width_meter}:{request.height_meter} meters)
- Design the background based on the theme "{request.message}"
- Analyze the theme and create appropriate visual elements, colors, and atmosphere
- For example:
  * Traffic safety → traffic signs, roads, safety symbols
  * Cultural festival → cultural patterns, festive elements
  * Environmental campaign → nature, green elements
  * Technology event → modern, tech patterns
- NO TEXT, NO WORDS, NO LETTERS, NO KOREAN CHARACTERS - only visual design elements
- Include decorative borders, gradients, patterns, and ornamental elements
- Leave the center area clean for text overlay
- Professional, suitable for Korean government institutions and public events

Important: DO NOT include any text, Korean characters, numbers, or words. Only create the background design that matches the theme."""

        response = client.models.generate_content(
            model="gemini-2.5-flash-image",
            contents=prompt,
            config=GenerateContentConfig(
                response_modalities=[Modality.IMAGE]
            )
        )

        # 배경 이미지 추출
        background_bytes = None
        for part in response.candidates[0].content.parts:
            if part.inline_data is not None:
                background_bytes = part.inline_data.data
                break

        if not background_bytes:
            raise HTTPException(status_code=500, detail="배경 이미지 생성 실패")

        # 배경 이미지를 PIL로 로드
        bg_image = Image.open(BytesIO(background_bytes))

        # 요청한 크기로 리사이즈
        bg_image = bg_image.resize((width_px, height_px), Image.LANCZOS)

        # 2단계: 텍스트 오버레이 (미리보기용)
        preview_image = bg_image.copy()
        draw = ImageDraw.Draw(preview_image)

        # 폰트 설정 시도 (시스템 폰트 사용)
        try:
            # 다양한 한글 폰트 경로 시도
            font_paths = [
                "/System/Library/Fonts/Supplemental/AppleGothic.ttf",  # macOS
                "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf",  # Linux
                "C:\\Windows\\Fonts\\malgun.ttf",  # Windows
            ]

            font_main = None
            font_small = None

            for font_path in font_paths:
                if os.path.exists(font_path):
                    font_main = ImageFont.truetype(font_path, size=int(height_px * 0.15))
                    font_small = ImageFont.truetype(font_path, size=int(height_px * 0.06))
                    break

            # 폰트를 찾지 못한 경우 기본 폰트 사용
            if not font_main:
                font_main = ImageFont.load_default()
                font_small = ImageFont.load_default()

        except Exception:
            font_main = ImageFont.load_default()
            font_small = ImageFont.load_default()

        # 주 문구 (중앙)
        bbox = draw.textbbox((0, 0), request.message, font=font_main)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]

        x = (width_px - text_width) / 2
        y = (height_px - text_height) / 2 - height_px * 0.1

        # 텍스트 그림자 효과
        shadow_offset = 4
        draw.text((x + shadow_offset, y + shadow_offset), request.message, fill=(0, 0, 0, 128), font=font_main)
        draw.text((x, y), request.message, fill='white', font=font_main)

        # 하단 정보
        bottom_y = height_px - height_px * 0.12

        # 날짜/시간 & 장소 (좌측)
        if request.date or request.location:
            info_parts = []
            if request.date:
                info_parts.append(request.date)
            if request.location:
                info_parts.append(request.location)
            info_text = " | ".join(info_parts)
            draw.text((50, bottom_y), info_text, fill='white', font=font_small)

        # 주최/주관 (우측)
        if request.organizer or request.host:
            org_parts = []
            if request.organizer:
                org_parts.append(f"주최: {request.organizer}")
            if request.host:
                org_parts.append(f"주관: {request.host}")
            org_text = " | ".join(org_parts)

            bbox = draw.textbbox((0, 0), org_text, font=font_small)
            org_width = bbox[2] - bbox[0]
            draw.text((width_px - org_width - 50, bottom_y), org_text, fill='white', font=font_small)

        # 3단계: 이미지를 Base64로 인코딩
        # 배경 이미지 (다운로드용)
        bg_buffer = BytesIO()
        bg_image.save(bg_buffer, format='PNG', quality=95)
        bg_base64 = base64.b64encode(bg_buffer.getvalue()).decode('utf-8')

        # 미리보기 이미지 (텍스트 포함)
        preview_buffer = BytesIO()
        preview_image.save(preview_buffer, format='PNG', quality=95)
        preview_base64 = base64.b64encode(preview_buffer.getvalue()).decode('utf-8')

        return {
            "success": True,
            "message": "현수막이 생성되었습니다",
            "background_image": bg_base64,  # 깨끗한 배경만 (다운로드용)
            "preview_image": preview_base64,  # 텍스트 포함 미리보기
            "filename": f"현수막배경_{request.width_meter}mx{request.height_meter}m.png"
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"현수막 생성 중 오류 발생: {str(e)}")


# ============================================
# PDF → PPT 변환
# ============================================

class PDFtoPPTRequest(BaseModel):
    pdf_file: str  # Base64 인코딩된 PDF
    slide_style: str = "professional"  # professional, modern, minimal


@app.post("/api/pdf-to-ppt")
async def pdf_to_ppt(request: PDFtoPPTRequest):
    """
    PDF 분석 후 PPT 생성
    - Gemini로 PDF 내용 분석 및 슬라이드 구조화
    - python-pptx로 PPT 파일 생성
    """
    try:
        # 1. PDF 디코딩
        pdf_bytes = base64.b64decode(request.pdf_file)

        # 2. Gemini로 PDF 분석
        prompt = """이 PDF 문서를 분석하여 프레젠테이션용 슬라이드 구조로 변환해주세요.

다음 JSON 형식으로 반드시 반환해주세요:
{
  "title": "프레젠테이션 제목",
  "slides": [
    {
      "type": "title",
      "title": "메인 제목",
      "subtitle": "부제목"
    },
    {
      "type": "content",
      "title": "슬라이드 제목",
      "bullets": ["포인트 1", "포인트 2", "포인트 3"]
    }
  ]
}

규칙:
- 각 슬라이드는 3-5개의 핵심 포인트만 포함
- 전체 10-15개 슬라이드로 구성
- 간결하고 명확하게 작성
- 첫 슬라이드는 반드시 type: "title"
- 나머지는 type: "content"
- JSON만 반환 (다른 텍스트 없이)"""

        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=[
                prompt,
                Part.from_bytes(data=pdf_bytes, mime_type="application/pdf")
            ]
        )

        # JSON 파싱
        response_text = response.text.strip()

        # ```json 제거
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]

        response_text = response_text.strip()

        try:
            slide_data = json.loads(response_text)
        except json.JSONDecodeError:
            raise HTTPException(status_code=500, detail="AI 응답을 파싱할 수 없습니다")

        # 3. PPT 생성
        prs = Presentation()

        # 슬라이드 크기 설정 (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        for slide_info in slide_data.get("slides", []):
            if slide_info.get("type") == "title":
                # 타이틀 슬라이드
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = slide_info.get("title", "")
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_info.get("subtitle", "")

            elif slide_info.get("type") == "content":
                # 내용 슬라이드
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_info.get("title", "")

                # 본문 텍스트
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()

                    for bullet in slide_info.get("bullets", []):
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(18)

        # 4. PPT를 바이트로 변환
        ppt_stream = BytesIO()
        prs.save(ppt_stream)
        ppt_bytes = ppt_stream.getvalue()
        ppt_base64 = base64.b64encode(ppt_bytes).decode('utf-8')

        return {
            "success": True,
            "message": "PPT가 생성되었습니다",
            "ppt_file": ppt_base64,
            "filename": f"{slide_data.get('title', 'presentation')}.pptx",
            "slide_count": len(slide_data.get("slides", []))
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF 변환 중 오류 발생: {str(e)}")


# AI 문제풀이 엔드포인트
class SolveProblemRequest(BaseModel):
    image: str  # base64 인코딩된 이미지
    subject: str = None  # 선택적 과목 정보

@app.post("/api/solve-problem")
async def solve_problem(request: SolveProblemRequest):
    """AI가 문제 이미지를 분석하고 해설을 제공합니다"""
    try:
        # base64 디코딩
        image_data = base64.b64decode(request.image)

        # MIME 타입 자동 감지
        mime_type = 'image/jpeg'
        if image_data[:8] == b'\x89PNG\r\n\x1a\n':
            mime_type = 'image/png'
        elif image_data[:2] == b'\xff\xd8':
            mime_type = 'image/jpeg'
        elif image_data[:6] in (b'GIF87a', b'GIF89a'):
            mime_type = 'image/gif'

        # 프롬프트 생성
        subject_info = f" (과목: {request.subject})" if request.subject else ""
        prompt = f"""다음 이미지에 있는 문제를 분석하고 상세한 해설을 제공해주세요{subject_info}.

**중요: 응답은 반드시 아래의 정확한 JSON 형식으로만 작성하세요. 다른 텍스트는 포함하지 마세요.**

{{
    "problem": "문제 내용 (원문 그대로)",
    "answer": "정답",
    "explanation": "상세한 해설 텍스트 (아래 지침 참고)",
    "concept": "핵심 개념과 학습 팁"
}}

상세 작성 지침:

1. **problem 필드**:
   - 문제 텍스트를 정확히 원문 그대로 작성
   - 지문이 있다면 지문도 포함
   - 보기(선택지)가 있다면 모두 포함

2. **answer 필드**:
   - 정답을 명확하고 간결하게 작성
   - 객관식이라면 번호와 내용 모두 표시 (예: "3번. extremely")

3. **explanation 필드** (가장 중요! - 반드시 문자열로 작성):
   영어 문제인 경우 다음 내용을 포함하여 하나의 긴 텍스트로 작성:

   📝 문제 번역
   (문제와 지문을 한국어로 번역)

   📋 보기 번역
   (모든 선택지를 한국어로 번역)

   📚 어휘 정리
   • 단어1: 뜻 (예문)
   • 단어2: 뜻 (예문)
   ...

   💡 풀이 과정
   1. 단계별 상세 설명
   2. ...

   ✅ 정답 근거
   (왜 이 답이 정답인지 명확히 설명)

   ❌ 오답 분석
   ① 왜 틀렸는지
   ② 왜 틀렸는지
   ...

   수학/과학 문제인 경우:
   - 문제 이해, 풀이 과정, 공식/원리, 주의사항을 상세히 작성

4. **concept 필드** (반드시 문자열로 작성):
   핵심 개념, 학습 팁, 관련 지식을 하나의 텍스트로 작성

**작성 원칙:**
- explanation과 concept는 반드시 문자열(string)로 작성
- 객체나 배열을 사용하지 말 것
- 초보자도 이해할 수 있도록 쉽고 친절하게 설명
- 모든 설명은 한국어로 작성
- 이모지를 사용해서 구조화하되, 하나의 긴 텍스트로 작성
- 충분히 상세하게 작성 (간략하게 하지 말 것!)

**JSON 응답 예시:**
{{
    "problem": "원문 문제...",
    "answer": "정답",
    "explanation": "📝 문제 번역\\n...\\n\\n📋 보기 번역\\n...\\n\\n📚 어휘 정리\\n• word: 뜻 (예문)\\n...",
    "concept": "🎯 핵심 개념\\n...\\n\\n📖 추가 설명\\n..."
}}"""

        # Google AI API 직접 사용 (앱과 동일)
        import google.generativeai as genai_sdk

        # API 키로 직접 연결
        genai_sdk.configure(api_key='AIzaSyD6_hNN5EyO_-OSmDrx3aFJzDwZ9xKMkTE')

        # 모델 초기화 (2.5 Pro - 최고 성능 모델)
        model = genai_sdk.GenerativeModel('gemini-2.5-pro')

        # 이미지를 PIL Image로 변환
        from PIL import Image
        import io
        image = Image.open(io.BytesIO(image_data))

        # API 호출
        response = model.generate_content(
            [prompt, image],
            generation_config=genai_sdk.GenerationConfig(
                temperature=0.4,
                max_output_tokens=8192,  # 더 긴 상세한 해설을 위해 증가
            )
        )

        # 응답 파싱
        result_text = response.text.strip()

        # 디버깅: 원본 응답 출력
        print(f"=== Gemini 2.5 Pro 원본 응답 (처음 500자) ===")
        print(result_text[:500])
        print("=" * 50)

        # JSON 추출 (코드 블록 제거)
        if '```json' in result_text:
            result_text = result_text.split('```json')[1].split('```')[0].strip()
        elif '```' in result_text:
            result_text = result_text.split('```')[1].split('```')[0].strip()

        try:
            result = json.loads(result_text)
        except json.JSONDecodeError:
            # JSON 파싱 실패 시 전체 텍스트를 explanation으로 반환
            print(f"JSON 파싱 실패. 전체 텍스트를 explanation으로 반환")
            return {
                "success": True,
                "problem": "문제를 인식했습니다.",
                "answer": "위 해설을 참고하세요.",
                "explanation": result_text,
                "concept": ""
            }

        # explanation과 concept가 객체인 경우 문자열로 변환
        explanation = result.get("explanation", "")
        if isinstance(explanation, dict):
            # 객체를 보기 좋은 문자열로 변환
            parts = []

            # 문제 번역
            if "문제 번역" in explanation:
                parts.append(f"📝 문제 번역\n{explanation['문제 번역']}\n")

            # 보기 번역
            if "보기 번역" in explanation:
                parts.append(f"📋 보기 번역\n{explanation['보기 번역']}\n")

            # 어휘 정리
            if "어휘 정리" in explanation:
                parts.append("📚 어휘 정리")
                for vocab in explanation["어휘 정리"]:
                    parts.append(f"\n• {vocab['word']}: {vocab['meaning']}")
                    if 'example' in vocab:
                        parts.append(f"  예문: {vocab['example']}")
                parts.append("\n")

            # 풀이 과정
            if "풀이 과정" in explanation:
                parts.append("💡 풀이 과정")
                for step in explanation["풀이 과정"]:
                    parts.append(f"\n{step}")
                parts.append("\n")

            # 정답 근거
            if "정답 근거" in explanation:
                parts.append(f"✅ 정답 근거\n{explanation['정답 근거']}\n")

            # 오답 분석
            if "오답 분석" in explanation:
                parts.append("❌ 오답 분석")
                for key, value in explanation["오답 분석"].items():
                    parts.append(f"\n{key} {value}")
                parts.append("\n")

            explanation = "\n".join(parts)

        concept = result.get("concept", "")
        if isinstance(concept, dict):
            # 객체를 보기 좋은 문자열로 변환
            parts = []
            if "핵심 개념" in concept:
                parts.append(f"🎯 핵심 개념\n{concept['핵심 개념']}\n")
            if "추가 설명" in concept:
                parts.append(f"📖 추가 설명\n{concept['추가 설명']}")
            concept = "\n".join(parts)

        return {
            "success": True,
            "problem": result.get("problem", ""),
            "answer": result.get("answer", ""),
            "explanation": explanation,
            "concept": concept
        }

    except json.JSONDecodeError as e:
        # JSON 파싱 실패 시 원본 텍스트 반환
        print(f"JSONDecodeError 발생: {str(e)}")
        return {
            "success": True,
            "problem": "문제를 인식했습니다.",
            "answer": "",
            "explanation": result_text if 'result_text' in locals() else "분석 중 오류가 발생했습니다.",
            "concept": ""
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"문제 분석 중 오류 발생: {str(e)}")


# ==================== AI 이미지 스타일 관리 API ====================

class AnalyzeStyleRequest(BaseModel):
    style_id: str
    reference_image: str  # Base64


class AnalyzeStyleOnlyRequest(BaseModel):
    reference_image: str  # Base64


class SaveStyleRequest(BaseModel):
    style_id: str
    style_prompt: str


class ConvertWithStylePromptRequest(BaseModel):
    image: str  # Base64
    style_prompt: str  # 저장된 스타일 프롬프트


@app.post("/api/analyze-style")
async def analyze_style(request: AnalyzeStyleRequest):
    """레퍼런스 이미지를 Gemini 2.5 Pro로 분석하여 스타일 프롬프트 생성 및 Firestore 저장"""
    try:
        # Base64 디코딩
        image_bytes = base64.b64decode(request.reference_image)

        # Gemini 2.5 Pro로 이미지 분석 (참고용 짧은 분석)
        analysis_prompt = """Analyze this image's visual style in 150 characters or less.

Focus on:
- Art style (anime, cartoon, painting, clay, etc.)
- Key visual characteristics
- Color tone
- Texture/material

Keep it brief and descriptive. This is just a reference for creating a detailed prompt later."""

        response = client.models.generate_content(
            model="gemini-2.5-pro",
            contents=[
                Part.from_bytes(data=image_bytes, mime_type="image/jpeg"),
                analysis_prompt
            ]
        )

        style_prompt = response.text

        # Firestore에 저장 (프롬프트만 저장, 이미지는 너무 커서 제외)
        style_data = {
            'id': request.style_id,
            'prompt': style_prompt,
            'created_at': firestore.SERVER_TIMESTAMP
        }

        firestore_db = get_firestore_client()
        firestore_db.collection('ai_image_styles').document(request.style_id).set(style_data)

        return {
            "success": True,
            "message": f"스타일 '{request.style_id}'이 성공적으로 분석되고 저장되었습니다.",
            "style_prompt": style_prompt
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"스타일 분석 중 오류 발생: {str(e)}")


@app.post("/api/analyze-style-only")
async def analyze_style_only(request: AnalyzeStyleOnlyRequest):
    """레퍼런스 이미지를 Gemini 2.5 Pro로 분석만 수행 (저장 안 함)"""
    try:
        # Base64 디코딩
        image_bytes = base64.b64decode(request.reference_image)

        # Gemini 2.5 Pro로 이미지 분석 (참고용 짧은 분석)
        analysis_prompt = """Analyze this image's visual style in 150 characters or less.

Focus on:
- Art style (anime, cartoon, painting, clay, etc.)
- Key visual characteristics
- Color tone
- Texture/material

Keep it brief and descriptive. This is just a reference for creating a detailed prompt later."""

        response = client.models.generate_content(
            model="gemini-2.5-pro",
            contents=[
                Part.from_bytes(data=image_bytes, mime_type="image/jpeg"),
                analysis_prompt
            ]
        )

        style_prompt = response.text

        return {
            "success": True,
            "message": "스타일 분석이 완료되었습니다. 프롬프트를 확인하고 수정한 후 저장하세요.",
            "style_prompt": style_prompt
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"스타일 분석 중 오류 발생: {str(e)}")


@app.post("/api/save-style")
async def save_style(request: SaveStyleRequest):
    """스타일 프롬프트를 Firestore에 저장"""
    try:
        # Firestore에 저장
        style_data = {
            'id': request.style_id,
            'prompt': request.style_prompt,
            'created_at': firestore.SERVER_TIMESTAMP
        }

        firestore_db = get_firestore_client()
        firestore_db.collection('ai_image_styles').document(request.style_id).set(style_data)

        return {
            "success": True,
            "message": f"스타일 '{request.style_id}'이 성공적으로 저장되었습니다."
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"스타일 저장 중 오류 발생: {str(e)}")


@app.get("/api/list-styles")
async def list_styles():
    """저장된 모든 스타일 목록 조회"""
    try:
        firestore_db = get_firestore_client()
        styles_ref = firestore_db.collection('ai_image_styles')
        docs = styles_ref.stream()

        styles = []
        for doc in docs:
            data = doc.to_dict()
            styles.append({
                'id': data.get('id'),
                'prompt': data.get('prompt'),
                'created_at': data.get('created_at')
            })

        # 최신순 정렬
        styles.sort(key=lambda x: x.get('created_at') or '', reverse=True)

        return {
            "success": True,
            "styles": styles
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"스타일 목록 조회 중 오류 발생: {str(e)}")


@app.post("/api/convert-with-style-prompt")
async def convert_with_style_prompt(request: ConvertWithStylePromptRequest):
    """저장된 스타일 프롬프트를 사용하여 이미지 변환"""
    try:
        # 사용자 이미지 디코딩
        user_image_bytes = base64.b64decode(request.image)

        # 스타일 프롬프트를 사용한 변환 프롬프트 생성
        conversion_prompt = f"""Transform this photograph to match the following artistic style:

{request.style_prompt}

CRITICAL REQUIREMENTS:
1. PRESERVE THE ORIGINAL COMPOSITION: Keep the person's exact pose, body position, facial expression, and overall layout exactly as shown in the photograph.
2. APPLY THE STYLE: Completely redraw/repaint the entire image using the artistic style described above.
3. The person's identity, pose, and expression must remain identical to the original.
4. The background should also be recreated in the same artistic style.
5. Create a cohesive, unified artwork where everything is rendered in the specified style.

Transform this image into a perfect recreation in the described artistic style while maintaining the original composition."""

        # Gemini 2.5 Flash Image로 변환
        response = client.models.generate_content(
            model="gemini-2.5-flash-image",
            contents=[
                Part.from_bytes(data=user_image_bytes, mime_type="image/jpeg"),
                conversion_prompt
            ],
            config=GenerateContentConfig(
                response_modalities=[Modality.IMAGE]
            )
        )

        # 응답 검증
        if not response or not response.candidates:
            raise HTTPException(
                status_code=500,
                detail="Gemini API에서 응답을 받지 못했습니다. 이미지가 정책을 위반했거나 서버 오류가 발생했을 수 있습니다."
            )

        # 안전 필터 확인
        candidate = response.candidates[0]
        if hasattr(candidate, 'finish_reason') and candidate.finish_reason:
            finish_reason = str(candidate.finish_reason)
            if 'SAFETY' in finish_reason:
                raise HTTPException(
                    status_code=400,
                    detail=f"안전 필터에 의해 차단되었습니다: {finish_reason}"
                )
            elif finish_reason not in ['STOP', 'FINISH_REASON_STOP']:
                raise HTTPException(
                    status_code=500,
                    detail=f"생성 중단됨: {finish_reason}"
                )

        if not candidate.content or not candidate.content.parts:
            # 디버깅을 위한 상세 정보
            error_detail = "생성된 콘텐츠가 없습니다."
            if hasattr(candidate, 'finish_reason'):
                error_detail += f" (finish_reason: {candidate.finish_reason})"
            raise HTTPException(
                status_code=500,
                detail=error_detail
            )

        # 생성된 이미지 추출
        for part in response.candidates[0].content.parts:
            if part.inline_data is not None:
                result_image_base64 = base64.b64encode(part.inline_data.data).decode('utf-8')
                return {
                    "success": True,
                    "message": "스타일 변환이 완료되었습니다",
                    "processed_image": result_image_base64
                }

        raise HTTPException(status_code=500, detail="이미지 생성에 실패했습니다")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"이미지 변환 중 오류 발생: {str(e)}")


class DeleteStyleRequest(BaseModel):
    style_id: str


@app.post("/api/delete-style")
async def delete_style(request: DeleteStyleRequest):
    """Firestore에서 스타일 삭제"""
    try:
        firestore_db = get_firestore_client()
        doc_ref = firestore_db.collection('ai_image_styles').document(request.style_id)

        # 문서가 존재하는지 확인
        if not doc_ref.get().exists:
            raise HTTPException(status_code=404, detail=f"스타일 '{request.style_id}'을 찾을 수 없습니다")

        # 문서 삭제
        doc_ref.delete()

        return {
            "success": True,
            "message": f"스타일 '{request.style_id}'이 성공적으로 삭제되었습니다."
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"스타일 삭제 중 오류 발생: {str(e)}")


# ========================================
# GemSem (교육 콘텐츠 제작) API
# ========================================

# ========================================
# 과목 관리 API
# ========================================

class SubjectCreateRequest(BaseModel):
    category: str  # "수능", "토익", "자격증"
    name: str  # 과목명
    description: str = ""

class SubjectUpdateRequest(BaseModel):
    name: str = None
    description: str = None

@app.post("/api/gemsem/subjects")
async def create_subject(request: SubjectCreateRequest):
    """과목 생성"""
    try:
        db = get_firestore_client()
        subject_ref = db.collection('gemsem_subjects').document()

        subject_data = {
            'id': subject_ref.id,
            'category': request.category,
            'name': request.name,
            'description': request.description,
            'created_at': firestore.SERVER_TIMESTAMP
        }

        subject_ref.set(subject_data)

        return {
            "success": True,
            "subject": subject_data,
            "message": f"과목 '{request.name}'이(가) 생성되었습니다."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"과목 생성 중 오류: {str(e)}")

@app.get("/api/gemsem/subjects")
async def list_subjects(category: str = None):
    """과목 목록 조회 (카테고리별 필터링 지원)"""
    try:
        db = get_firestore_client()
        subjects_ref = db.collection('gemsem_subjects')

        if category:
            subjects_ref = subjects_ref.where('category', '==', category)

        subjects = []
        for doc in subjects_ref.stream():
            subject = doc.to_dict()
            subjects.append(subject)

        return {
            "success": True,
            "subjects": subjects,
            "count": len(subjects)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"과목 조회 중 오류: {str(e)}")

@app.put("/api/gemsem/subjects/{subject_id}")
async def update_subject(subject_id: str, request: SubjectUpdateRequest):
    """과목 수정"""
    try:
        db = get_firestore_client()
        subject_ref = db.collection('gemsem_subjects').document(subject_id)
        subject_doc = subject_ref.get()

        if not subject_doc.exists:
            raise HTTPException(status_code=404, detail="과목을 찾을 수 없습니다.")

        update_data = {}
        if request.name is not None:
            update_data['name'] = request.name
        if request.description is not None:
            update_data['description'] = request.description

        if update_data:
            update_data['updated_at'] = firestore.SERVER_TIMESTAMP
            subject_ref.update(update_data)

        return {
            "success": True,
            "message": "과목이 수정되었습니다."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"과목 수정 중 오류: {str(e)}")

@app.delete("/api/gemsem/subjects/{subject_id}")
async def delete_subject(subject_id: str):
    """과목 삭제"""
    try:
        db = get_firestore_client()
        subject_ref = db.collection('gemsem_subjects').document(subject_id)
        subject_doc = subject_ref.get()

        if not subject_doc.exists:
            raise HTTPException(status_code=404, detail="과목을 찾을 수 없습니다.")

        subject_ref.delete()

        return {
            "success": True,
            "message": "과목이 삭제되었습니다."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"과목 삭제 중 오류: {str(e)}")

# ========================================
# 교과서/문제집 API
# ========================================

class TextbookUploadRequest(BaseModel):
    title: str
    category: str  # "수능", "토익", "자격증"
    subject: str  # 과목명
    publisher: str
    pdf_base64: str  # Base64 encoded PDF

class TextbookAnalysisRequest(BaseModel):
    textbook_id: str

class ConceptExtractionResponse(BaseModel):
    success: bool
    textbook_id: str
    concepts: list
    message: str

@app.post("/api/gemsem/textbooks/upload")
async def upload_textbook(request: TextbookUploadRequest):
    """교과서 PDF 업로드 및 기본 정보 저장"""
    try:
        # Base64 PDF 디코딩
        pdf_bytes = base64.b64decode(request.pdf_base64)

        # Firestore에 교과서 정보 저장
        db = get_firestore_client()
        textbook_ref = db.collection('gemsem_textbooks').document()

        textbook_data = {
            'id': textbook_ref.id,
            'title': request.title,
            'category': request.category,  # "수능", "토익", "자격증"
            'subject': request.subject,
            'publisher': request.publisher,
            'pdf_base64': request.pdf_base64,  # PDF 저장
            'pdf_size_bytes': len(pdf_bytes),
            'upload_date': firestore.SERVER_TIMESTAMP,
            'status': 'uploaded',  # uploaded → processing → completed
            'processing_progress': 0
        }

        textbook_ref.set(textbook_data)

        return {
            "success": True,
            "textbook_id": textbook_ref.id,
            "message": f"교과서 '{request.title}'이 성공적으로 업로드되었습니다.",
            "pdf_size_mb": round(len(pdf_bytes) / 1024 / 1024, 2)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"교과서 업로드 중 오류: {str(e)}")


@app.post("/api/gemsem/textbooks/{textbook_id}/analyze")
async def analyze_textbook(textbook_id: str):
    """교과서 PDF에서 개념 추출 (Gemini 2.0 Flash 사용)"""
    try:
        db = get_firestore_client()
        textbook_ref = db.collection('gemsem_textbooks').document(textbook_id)
        textbook_doc = textbook_ref.get()

        if not textbook_doc.exists:
            raise HTTPException(status_code=404, detail="교과서를 찾을 수 없습니다.")

        textbook_data = textbook_doc.to_dict()

        # 상태 업데이트: processing
        textbook_ref.update({
            'status': 'processing',
            'processing_start_time': firestore.SERVER_TIMESTAMP
        })

        # PDF를 Gemini 2.0 Flash로 직접 분석
        pdf_base64 = textbook_data.get('pdf_base64')
        if not pdf_base64:
            raise HTTPException(status_code=400, detail="PDF 파일이 없습니다.")

        concept_prompt = f"""이 PDF는 {textbook_data['grade_level']} {textbook_data['subject']} 교과서입니다.
PDF 전체 내용을 분석하여 주요 교육 개념들을 추출해주세요.

각 개념에 대해 다음 정보를 포함하세요:
- name: 개념의 이름
- description: 개념에 대한 설명 (2-3문장)
- formula: 수식이 있다면 포함 (없으면 빈 문자열)
- difficulty: 난이도 (하/중/상)
- chapter: 어느 단원/챕터에 속하는지
- prerequisites: 선수 개념 목록 (배열)

JSON 형식으로만 출력하세요:
{{
  "concepts": [
    {{
      "name": "개념명",
      "description": "개념 설명",
      "formula": "수식 (선택사항)",
      "difficulty": "중",
      "chapter": "단원명",
      "prerequisites": ["선수개념1", "선수개념2"]
    }}
  ]
}}

최소 5개 이상의 주요 개념을 추출해주세요."""

        # Gemini에 PDF 전송
        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=[
                {
                    "role": "user",
                    "parts": [
                        {"text": concept_prompt},
                        {
                            "inline_data": {
                                "mime_type": "application/pdf",
                                "data": pdf_base64
                            }
                        }
                    ]
                }
            ]
        )

        # JSON 파싱
        concepts_text = response.text.strip()
        # Markdown 코드 블록 제거
        if concepts_text.startswith("```json"):
            concepts_text = concepts_text[7:]
        if concepts_text.startswith("```"):
            concepts_text = concepts_text[3:]
        if concepts_text.endswith("```"):
            concepts_text = concepts_text[:-3]

        concepts_data = json.loads(concepts_text.strip())

        # Firestore에 개념 저장
        concepts_collection = db.collection('gemsem_concepts')
        saved_concepts = []

        for concept in concepts_data.get('concepts', []):
            concept_ref = concepts_collection.document()
            concept_doc = {
                'id': concept_ref.id,
                'textbook_id': textbook_id,
                'name': concept.get('name'),
                'description': concept.get('description'),
                'formula': concept.get('formula', ''),
                'difficulty': concept.get('difficulty', '중'),
                'chapter': concept.get('chapter', ''),
                'prerequisites': concept.get('prerequisites', []),
                'created_at': firestore.SERVER_TIMESTAMP,
                'verified': False  # 검수 전
            }
            concept_ref.set(concept_doc)
            saved_concepts.append(concept_doc)

        # 교과서 상태 업데이트: completed
        textbook_ref.update({
            'status': 'completed',
            'processing_end_time': firestore.SERVER_TIMESTAMP,
            'concepts_count': len(saved_concepts),
            'processing_progress': 100
        })

        return {
            "success": True,
            "textbook_id": textbook_id,
            "concepts": saved_concepts,
            "message": f"{len(saved_concepts)}개의 개념이 추출되었습니다."
        }

    except json.JSONDecodeError as e:
        textbook_ref.update({'status': 'failed', 'error_message': f"JSON 파싱 오류: {str(e)}"})
        raise HTTPException(status_code=500, detail=f"개념 추출 결과를 파싱할 수 없습니다: {str(e)}")
    except Exception as e:
        textbook_ref.update({'status': 'failed', 'error_message': str(e)})
        raise HTTPException(status_code=500, detail=f"교과서 분석 중 오류: {str(e)}")


@app.get("/api/gemsem/textbooks")
async def list_textbooks(subject: str = None):
    """저장된 교과서 목록 조회 (과목별 필터링 지원)"""
    try:
        db = get_firestore_client()
        textbooks_ref = db.collection('gemsem_textbooks')

        # 과목 필터링
        if subject:
            textbooks_ref = textbooks_ref.where('subject', '==', subject)

        textbooks_ref = textbooks_ref.order_by('upload_date', direction=firestore.Query.DESCENDING).limit(50)
        textbooks = []

        for doc in textbooks_ref.stream():
            textbook = doc.to_dict()
            textbooks.append(textbook)

        return {
            "success": True,
            "textbooks": textbooks,
            "count": len(textbooks)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"교과서 목록 조회 중 오류: {str(e)}")


@app.get("/api/gemsem/textbooks/{textbook_id}/concepts")
async def get_textbook_concepts(textbook_id: str):
    """특정 교과서의 개념 목록 조회"""
    try:
        db = get_firestore_client()
        concepts_ref = db.collection('gemsem_concepts').where('textbook_id', '==', textbook_id)
        concepts = []

        for doc in concepts_ref.stream():
            concept = doc.to_dict()
            concepts.append(concept)

        return {
            "success": True,
            "textbook_id": textbook_id,
            "concepts": concepts,
            "count": len(concepts)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"개념 조회 중 오류: {str(e)}")


@app.put("/api/gemsem/concepts/{concept_id}/verify")
async def verify_concept(concept_id: str):
    """개념 검수 완료 처리"""
    try:
        db = get_firestore_client()
        concept_ref = db.collection('gemsem_concepts').document(concept_id)
        concept_doc = concept_ref.get()

        if not concept_doc.exists:
            raise HTTPException(status_code=404, detail="개념을 찾을 수 없습니다.")

        concept_ref.update({
            'verified': True,
            'verified_at': firestore.SERVER_TIMESTAMP
        })

        return {
            "success": True,
            "message": "개념이 검수 완료되었습니다."
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"개념 검수 중 오류: {str(e)}")


@app.get("/api/gemsem/concepts/all")
async def get_all_concepts(subject: str = None, textbook_id: str = None):
    """전체 개념 조회 (과목별/교과서별 필터링 지원)"""
    try:
        db = get_firestore_client()
        concepts_ref = db.collection('gemsem_concepts')

        # 필터링 조건 적용
        if textbook_id:
            concepts_ref = concepts_ref.where('textbook_id', '==', textbook_id)

        concepts = []
        for doc in concepts_ref.stream():
            concept = doc.to_dict()

            # 과목 필터링 (교과서 정보에서 가져와야 함)
            if subject:
                textbook_ref = db.collection('gemsem_textbooks').document(concept['textbook_id'])
                textbook_doc = textbook_ref.get()
                if textbook_doc.exists and textbook_doc.to_dict().get('subject') == subject:
                    # 교과서 정보 추가
                    concept['textbook_title'] = textbook_doc.to_dict().get('title')
                    concept['subject'] = textbook_doc.to_dict().get('subject')
                    concepts.append(concept)
            else:
                # 과목 필터가 없으면 교과서 정보만 추가
                textbook_ref = db.collection('gemsem_textbooks').document(concept['textbook_id'])
                textbook_doc = textbook_ref.get()
                if textbook_doc.exists:
                    concept['textbook_title'] = textbook_doc.to_dict().get('title')
                    concept['subject'] = textbook_doc.to_dict().get('subject')
                concepts.append(concept)

        return {
            "success": True,
            "concepts": concepts,
            "count": len(concepts)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"개념 조회 중 오류: {str(e)}")


# ========================================
# 문제집 분석 API
# ========================================

class WorkbookUploadRequest(BaseModel):
    title: str
    category: str  # "수능", "토익", "자격증"
    subject: str
    publisher: str
    pdf_base64: str

@app.post("/api/gemsem/workbooks/upload")
async def upload_workbook(request: WorkbookUploadRequest):
    """문제집 PDF 업로드 및 기본 정보 저장"""
    try:
        pdf_bytes = base64.b64decode(request.pdf_base64)

        db = get_firestore_client()
        workbook_ref = db.collection('gemsem_workbooks').document()

        workbook_data = {
            'id': workbook_ref.id,
            'title': request.title,
            'category': request.category,  # "수능", "토익", "자격증"
            'subject': request.subject,
            'publisher': request.publisher,
            'pdf_base64': request.pdf_base64,  # PDF 저장
            'pdf_size_bytes': len(pdf_bytes),
            'upload_date': firestore.SERVER_TIMESTAMP,
            'status': 'uploaded',
            'processing_progress': 0
        }

        workbook_ref.set(workbook_data)

        return {
            "success": True,
            "workbook_id": workbook_ref.id,
            "message": f"문제집 '{request.title}'이 성공적으로 업로드되었습니다.",
            "pdf_size_mb": round(len(pdf_bytes) / 1024 / 1024, 2)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"문제집 업로드 중 오류: {str(e)}")


@app.post("/api/gemsem/workbooks/{workbook_id}/analyze")
async def analyze_workbook(workbook_id: str):
    """문제집 PDF에서 문제 추출 (Gemini 2.0 Flash 사용)"""
    try:
        db = get_firestore_client()
        workbook_ref = db.collection('gemsem_workbooks').document(workbook_id)
        workbook_doc = workbook_ref.get()

        if not workbook_doc.exists:
            raise HTTPException(status_code=404, detail="문제집을 찾을 수 없습니다.")

        workbook_data = workbook_doc.to_dict()

        workbook_ref.update({
            'status': 'processing',
            'processing_start_time': firestore.SERVER_TIMESTAMP
        })

        # PDF를 Gemini 2.0 Flash로 직접 분석
        pdf_base64 = workbook_data.get('pdf_base64')
        if not pdf_base64:
            raise HTTPException(status_code=400, detail="PDF 파일이 없습니다.")

        problem_prompt = f"""이 PDF는 {workbook_data['grade_level']} {workbook_data['subject']} 문제집입니다.
PDF의 모든 문제를 추출하여 구조화된 데이터로 변환해주세요.

각 문제에 대해 다음 정보를 포함하세요:
- problem_number: 문제 번호 (문자열)
- problem_text: 문제 내용 전체
- problem_type: 문제 유형 (객관식/주관식/서술형)
- choices: 객관식인 경우 선택지 배열 (주관식이면 빈 배열)
- answer: 정답 (있으면 포함, 없으면 빈 문자열)
- difficulty: 난이도 추정 (하/중/상)
- concepts: 관련 개념 목록 (배열)
- page: 페이지 번호 (추정)

JSON 형식으로만 출력하세요:
{{
  "problems": [
    {{
      "problem_number": "1",
      "problem_text": "문제 내용",
      "problem_type": "객관식",
      "choices": ["선택지1", "선택지2", "선택지3", "선택지4"],
      "answer": "①",
      "difficulty": "중",
      "concepts": ["개념1", "개념2"],
      "page": 1
    }}
  ]
}}

최소 10개 이상의 문제를 추출해주세요."""

        # Gemini에 PDF 전송
        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=[
                {
                    "role": "user",
                    "parts": [
                        {"text": problem_prompt},
                        {
                            "inline_data": {
                                "mime_type": "application/pdf",
                                "data": pdf_base64
                            }
                        }
                    ]
                }
            ]
        )

        # JSON 파싱
        problems_text = response.text.strip()
        if problems_text.startswith("```json"):
            problems_text = problems_text[7:]
        if problems_text.startswith("```"):
            problems_text = problems_text[3:]
        if problems_text.endswith("```"):
            problems_text = problems_text[:-3]

        problems_data = json.loads(problems_text.strip())

        # Firestore에 문제 저장
        problems_collection = db.collection('gemsem_problems')
        saved_problems = []

        for problem in problems_data.get('problems', []):
            problem_ref = problems_collection.document()
            problem_id = problem_ref.id

            # 원본 문제는 자기 자신이 문제 원형(template)이 됨
            problem_doc = {
                'id': problem_id,
                'workbook_id': workbook_id,
                'category': workbook_data['category'],  # 섹션 정보 (수능/토익/자격증)
                'subject': workbook_data['subject'],  # 과목 정보
                'problem_number': problem.get('problem_number'),
                'problem_text': problem.get('problem_text'),
                'problem_type': problem.get('problem_type', '객관식'),
                'choices': problem.get('choices', []),
                'answer': problem.get('answer', ''),
                'difficulty': problem.get('difficulty', '중'),
                'concepts': problem.get('concepts', []),
                'page': problem.get('page', 1),
                'content_source': 'original',  # 저작권 주의!
                'template_id': problem_id,  # 원본 문제는 자기 자신이 원형
                'is_template': True,  # 이 문제는 원형임
                'derived_count': 0,  # 이 원형에서 파생된 문제 개수
                'created_at': firestore.SERVER_TIMESTAMP,
                'verified': False
            }
            problem_ref.set(problem_doc)
            saved_problems.append(problem_doc)

        workbook_ref.update({
            'status': 'completed',
            'processing_end_time': firestore.SERVER_TIMESTAMP,
            'problems_count': len(saved_problems),
            'processing_progress': 100
        })

        return {
            "success": True,
            "workbook_id": workbook_id,
            "problems": saved_problems,
            "message": f"{len(saved_problems)}개의 문제가 추출되었습니다."
        }

    except json.JSONDecodeError as e:
        workbook_ref.update({'status': 'failed', 'error_message': f"JSON 파싱 오류: {str(e)}"})
        raise HTTPException(status_code=500, detail=f"문제 추출 결과를 파싱할 수 없습니다: {str(e)}")
    except Exception as e:
        workbook_ref.update({'status': 'failed', 'error_message': str(e)})
        raise HTTPException(status_code=500, detail=f"문제집 분석 중 오류: {str(e)}")


@app.get("/api/gemsem/workbooks")
async def list_workbooks(subject: str = None):
    """저장된 문제집 목록 조회 (과목별 필터링 지원)"""
    try:
        db = get_firestore_client()
        workbooks_ref = db.collection('gemsem_workbooks')

        # 과목 필터링
        if subject:
            workbooks_ref = workbooks_ref.where('subject', '==', subject)

        workbooks_ref = workbooks_ref.order_by('upload_date', direction=firestore.Query.DESCENDING).limit(50)
        workbooks = []

        for doc in workbooks_ref.stream():
            workbook = doc.to_dict()
            workbooks.append(workbook)

        return {
            "success": True,
            "workbooks": workbooks,
            "count": len(workbooks)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"문제집 목록 조회 중 오류: {str(e)}")


@app.get("/api/gemsem/workbooks/{workbook_id}/problems")
async def get_workbook_problems(workbook_id: str):
    """특정 문제집의 문제 목록 조회"""
    try:
        db = get_firestore_client()
        problems_ref = db.collection('gemsem_problems').where('workbook_id', '==', workbook_id)
        problems = []

        for doc in problems_ref.stream():
            problem = doc.to_dict()
            problems.append(problem)

        return {
            "success": True,
            "workbook_id": workbook_id,
            "problems": problems,
            "count": len(problems)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"문제 조회 중 오류: {str(e)}")


@app.get("/api/gemsem/problems/all")
async def get_all_problems(
    subject: str = None,
    workbook_id: str = None,
    template_id: str = None,
    is_template: bool = None
):
    """전체 문제 조회 (과목별/문제집별/원형별 필터링 지원)"""
    try:
        db = get_firestore_client()
        problems_ref = db.collection('gemsem_problems')

        # 필터링 조건 적용
        if workbook_id:
            problems_ref = problems_ref.where('workbook_id', '==', workbook_id)
        if template_id:
            problems_ref = problems_ref.where('template_id', '==', template_id)
        if is_template is not None:
            problems_ref = problems_ref.where('is_template', '==', is_template)
        if subject:
            problems_ref = problems_ref.where('subject', '==', subject)

        problems = []
        for doc in problems_ref.stream():
            problem = doc.to_dict()
            problems.append(problem)

        return {
            "success": True,
            "problems": problems,
            "count": len(problems)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"문제 조회 중 오류: {str(e)}")


@app.get("/api/gemsem/problems/templates")
async def get_problem_templates(subject: str = None):
    """문제 원형 목록 조회 (파생 문제 개수 포함)"""
    try:
        db = get_firestore_client()
        problems_ref = db.collection('gemsem_problems').where('is_template', '==', True)

        # 과목 필터링
        if subject:
            problems_ref = problems_ref.where('subject', '==', subject)

        templates = []
        for doc in problems_ref.stream():
            template = doc.to_dict()
            templates.append(template)

        return {
            "success": True,
            "templates": templates,
            "count": len(templates)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"문제 원형 조회 중 오류: {str(e)}")


# ========================================
# 유사 문제 생성 API (수동 버튼)
# ========================================

class GenerateSimilarRequest(BaseModel):
    problem_id: str
    count: int = 3  # 기본 3개 생성

@app.post("/api/gemsem/problems/{problem_id}/generate-similar")
async def generate_similar_problems(problem_id: str, request: GenerateSimilarRequest):
    """원본 문제를 기반으로 유사 문제 생성 (Gemini 2.0 Flash)"""
    try:
        db = get_firestore_client()
        problem_ref = db.collection('gemsem_problems').document(problem_id)
        problem_doc = problem_ref.get()

        if not problem_doc.exists:
            raise HTTPException(status_code=404, detail="문제를 찾을 수 없습니다.")

        original_problem = problem_doc.to_dict()

        # Gemini 2.0 Flash로 유사 문제 생성
        similar_prompt = f"""다음 문제를 참고하여 유사한 문제를 {request.count}개 생성하세요.
- 개념은 동일하게 유지
- 난이도는 동일하게 유지
- 숫자와 상황만 변경
- 정답도 함께 제공

원본 문제:
{original_problem['problem_text']}

{f"보기: {original_problem.get('choices', [])}" if original_problem.get('choices') else ""}

출력 형식 (JSON):
{{
  "similar_problems": [
    {{
      "problem_text": "...",
      "choices": ["...", "...", "...", "..."],
      "answer": "...",
      "variation_type": "숫자 변경"
    }},
    ...
  ]
}}
"""

        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=similar_prompt
        )

        # JSON 파싱
        similar_text = response.text.strip()
        if similar_text.startswith("```json"):
            similar_text = similar_text[7:]
        if similar_text.startswith("```"):
            similar_text = similar_text[3:]
        if similar_text.endswith("```"):
            similar_text = similar_text[:-3]

        similar_data = json.loads(similar_text.strip())

        # Firestore에 유사 문제 저장
        problems_collection = db.collection('gemsem_problems')
        saved_similar = []

        # 원형 문제의 template_id 가져오기 (원본 문제가 원형이면 자기 자신의 ID)
        template_id = original_problem.get('template_id', problem_id)

        for idx, similar_prob in enumerate(similar_data.get('similar_problems', []), 1):
            similar_ref = problems_collection.document()
            similar_doc = {
                'id': similar_ref.id,
                'workbook_id': original_problem.get('workbook_id'),
                'category': original_problem.get('category'),  # 섹션 정보 유지 (수능/토익/자격증)
                'subject': original_problem.get('subject'),  # 과목 정보 유지
                'problem_number': f"{original_problem['problem_number']}-{idx}",  # 파생 문제 번호
                'problem_text': similar_prob.get('problem_text'),
                'problem_type': original_problem.get('problem_type'),
                'choices': similar_prob.get('choices', []),
                'answer': similar_prob.get('answer', ''),
                'difficulty': original_problem.get('difficulty'),
                'concepts': original_problem.get('concepts', []),
                'content_source': 'ai_generated',  # AI 생성 (저작권 Free)
                'template_id': template_id,  # 속한 문제 원형 ID
                'is_template': False,  # 파생 문제는 원형이 아님
                'source_problem_id': problem_id,  # 직접적인 소스 문제 ID
                'ai_generation_info': {
                    'model': 'gemini-2.0-flash-exp',
                    'variation_type': similar_prob.get('variation_type', '숫자 변경'),
                    'generated_from': problem_id
                },
                'created_at': firestore.SERVER_TIMESTAMP,
                'verified': False
            }
            similar_ref.set(similar_doc)
            saved_similar.append(similar_doc)

        # 원형 문제의 파생 카운트 업데이트
        template_ref = db.collection('gemsem_problems').document(template_id)
        template_doc = template_ref.get()
        if template_doc.exists:
            current_count = template_doc.to_dict().get('derived_count', 0)
            template_ref.update({'derived_count': current_count + len(saved_similar)})

        return {
            "success": True,
            "original_problem_id": problem_id,
            "similar_problems": saved_similar,
            "count": len(saved_similar),
            "message": f"{len(saved_similar)}개의 유사 문제가 생성되었습니다."
        }

    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"유사 문제 생성 결과를 파싱할 수 없습니다: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"유사 문제 생성 중 오류: {str(e)}")


# ========================================
# 해설 자동 생성 API
# ========================================

@app.post("/api/gemsem/problems/{problem_id}/generate-solution")
async def generate_solution(problem_id: str):
    """문제에 대한 단계별 해설 생성 (Gemini 2.0 Flash)"""
    try:
        db = get_firestore_client()
        problem_ref = db.collection('gemsem_problems').document(problem_id)
        problem_doc = problem_ref.get()

        if not problem_doc.exists:
            raise HTTPException(status_code=404, detail="문제를 찾을 수 없습니다.")

        problem = problem_doc.to_dict()

        # Gemini 2.0 Flash로 해설 생성
        solution_prompt = f"""다음 문제의 단계별 풀이를 작성하세요.

문제:
{problem['problem_text']}

{f"보기: {problem.get('choices', [])}" if problem.get('choices') else ""}

정답: {problem.get('answer', '')}

풀이를 다음 형식으로 작성하세요:

**Step 1: 문제 분석**
- 무엇을 구해야 하는지 파악
- 주어진 조건 정리

**Step 2: 개념 적용**
- 어떤 개념/공식을 사용할지 결정
- 공식 명시

**Step 3: 계산 과정**
- 단계별 계산 과정
- 중간 결과 표시

**Step 4: 답 확인**
- 최종 답 도출
- 답의 타당성 검증

각 단계를 명확하고 자세하게 설명하세요.
"""

        response = client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=solution_prompt
        )

        solution_text = response.text.strip()

        # Firestore에 해설 저장
        solutions_collection = db.collection('gemsem_solutions')
        solution_ref = solutions_collection.document()

        solution_doc = {
            'id': solution_ref.id,
            'problem_id': problem_id,
            'solution_text': solution_text,
            'solution_type': 'step_by_step',
            'created_by': 'ai',
            'model': 'gemini-2.0-flash-exp',
            'created_at': firestore.SERVER_TIMESTAMP,
            'verified': False
        }

        solution_ref.set(solution_doc)

        return {
            "success": True,
            "problem_id": problem_id,
            "solution": solution_doc,
            "message": "해설이 생성되었습니다."
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"해설 생성 중 오류: {str(e)}")


@app.get("/api/gemsem/problems/{problem_id}/solution")
async def get_solution(problem_id: str):
    """문제의 해설 조회"""
    try:
        db = get_firestore_client()
        solutions_ref = db.collection('gemsem_solutions').where('problem_id', '==', problem_id).limit(1)

        for doc in solutions_ref.stream():
            solution = doc.to_dict()
            return {
                "success": True,
                "solution": solution
            }

        return {
            "success": False,
            "message": "해설이 없습니다."
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"해설 조회 중 오류: {str(e)}")


# ============================================
# Bible App Analytics API
# ============================================

class BibleAnalyticsData(BaseModel):
    userId: str
    lastActiveDate: str
    totalDaysActive: int
    totalSessions: int
    firstInstallDate: str
    dailyActivityLog: dict
    versesSaved: int
    chaptersRead: int
    totalReadingTimeSeconds: int
    appVersion: str = "1.0.0"
    platform: str = "flutter"
    syncedAt: str = None

@app.post("/api/bible-analytics")
async def receive_bible_analytics(data: BibleAnalyticsData):
    """
    Bible 앱에서 보낸 analytics 데이터를 저장
    """
    try:
        # Firestore에 저장 - gogwan 프로젝트의 별도 컬렉션에 저장
        db_client = get_firestore_client()
        analytics_ref = db_client.collection('bible_analytics').document(data.userId)

        analytics_ref.set({
            'userId': data.userId,
            'lastActiveDate': data.lastActiveDate,
            'totalDaysActive': data.totalDaysActive,
            'totalSessions': data.totalSessions,
            'firstInstallDate': data.firstInstallDate,
            'dailyActivityLog': data.dailyActivityLog,
            'versesSaved': data.versesSaved,
            'chaptersRead': data.chaptersRead,
            'totalReadingTimeSeconds': data.totalReadingTimeSeconds,
            'appVersion': data.appVersion,
            'platform': data.platform,
            'lastSyncedAt': firestore.SERVER_TIMESTAMP
        }, merge=True)

        return {
            "success": True,
            "message": "Analytics data saved successfully",
            "userId": data.userId
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving analytics: {str(e)}")

@app.get("/api/bible-analytics/stats")
async def get_bible_analytics_stats():
    """
    Bible 앱의 전체 통계 조회
    """
    try:
        # gogwan 프로젝트의 Firestore 사용
        db_client = get_firestore_client()
        analytics_ref = db_client.collection('bible_analytics')

        # 전체 사용자 데이터 가져오기
        users = []
        for doc in analytics_ref.stream():
            user_data = doc.to_dict()
            users.append(user_data)

        if not users:
            return {
                "totalUsers": 0,
                "dau": 0,
                "mau": 0,
                "avgStreak": 0,
                "totalVersesSaved": 0,
                "totalChaptersRead": 0,
                "avgReadingTime": 0,
                "avgEngagement": 0
            }

        # 통계 계산
        from datetime import datetime, timedelta
        today = datetime.now().date()
        yesterday = today - timedelta(days=1)
        last_30_days = today - timedelta(days=30)

        dau = 0  # Daily Active Users
        mau = 0  # Monthly Active Users
        total_verses_saved = 0
        total_chapters_read = 0
        total_reading_time = 0

        for user in users:
            # DAU 계산 (오늘 활동한 사용자)
            last_active = datetime.fromisoformat(user.get('lastActiveDate', '')).date()
            if last_active == today or last_active == yesterday:
                dau += 1

            # MAU 계산 (최근 30일 활동)
            if last_active >= last_30_days:
                mau += 1

            total_verses_saved += user.get('versesSaved', 0)
            total_chapters_read += user.get('chaptersRead', 0)
            total_reading_time += user.get('totalReadingTimeSeconds', 0)

        avg_reading_time = total_reading_time / len(users) / 60 if users else 0  # minutes

        return {
            "totalUsers": len(users),
            "dau": dau,
            "mau": mau,
            "avgStreak": 0,  # TODO: Calculate average streak
            "totalVersesSaved": total_verses_saved,
            "totalChaptersRead": total_chapters_read,
            "avgReadingTime": round(avg_reading_time, 1),
            "avgEngagement": 0,  # TODO: Calculate engagement
            "users": users[:100]  # Return top 100 users
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching analytics: {str(e)}")

# ===== 게시판 API =====
class BoardPost(BaseModel):
    title: str
    content: str
    author: str
    password: str
    category: str = "자유게시판"

class BoardPostUpdate(BaseModel):
    title: str
    content: str
    password: str

class BoardComment(BaseModel):
    content: str
    author: str
    password: str

class BoardCommentUpdate(BaseModel):
    content: str
    password: str

class PasswordVerify(BaseModel):
    password: str

@app.get("/api/board/posts")
async def get_board_posts(category: str = "자유게시판", limit: int = 50, offset: int = 0):
    """
    게시판 글 목록 조회
    """
    try:
        db_client = get_firestore_client()
        posts_ref = db_client.collection('board_posts')

        # 카테고리 필터링 및 최신순 정렬
        query = posts_ref.where('category', '==', category).order_by('createdAt', direction=firestore.Query.DESCENDING).limit(limit)

        posts = []
        for doc in query.stream():
            post_data = doc.to_dict()
            post_data['id'] = doc.id
            posts.append(post_data)

        return {
            "success": True,
            "posts": posts,
            "total": len(posts)
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching posts: {str(e)}")

@app.get("/api/board/posts/{post_id}")
async def get_board_post(post_id: str):
    """
    게시글 상세 조회
    """
    try:
        db_client = get_firestore_client()
        post_ref = db_client.collection('board_posts').document(post_id)
        post_doc = post_ref.get()

        if not post_doc.exists:
            raise HTTPException(status_code=404, detail="Post not found")

        post_data = post_doc.to_dict()
        post_data['id'] = post_doc.id

        # 조회수 증가
        post_ref.update({
            'views': firestore.Increment(1)
        })
        post_data['views'] = post_data.get('views', 0) + 1

        return {
            "success": True,
            "post": post_data
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching post: {str(e)}")

@app.post("/api/board/posts")
async def create_board_post(post: BoardPost):
    """
    게시글 작성
    """
    try:
        db_client = get_firestore_client()
        posts_ref = db_client.collection('board_posts')

        # 비밀번호 해싱 (간단한 해시 사용)
        import hashlib
        password_hash = hashlib.sha256(post.password.encode()).hexdigest()

        post_data = {
            'title': post.title,
            'content': post.content,
            'author': post.author,
            'password': password_hash,
            'category': post.category,
            'createdAt': firestore.SERVER_TIMESTAMP,
            'updatedAt': firestore.SERVER_TIMESTAMP,
            'views': 0,
            'likes': 0
        }

        doc_ref = posts_ref.add(post_data)
        post_id = doc_ref[1].id

        return {
            "success": True,
            "message": "Post created successfully",
            "postId": post_id
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating post: {str(e)}")

@app.put("/api/board/posts/{post_id}")
async def update_board_post(post_id: str, post: BoardPostUpdate):
    """
    게시글 수정
    """
    try:
        import hashlib
        db_client = get_firestore_client()
        post_ref = db_client.collection('board_posts').document(post_id)
        post_doc = post_ref.get()

        if not post_doc.exists:
            raise HTTPException(status_code=404, detail="Post not found")

        # 비밀번호 검증
        stored_password = post_doc.to_dict().get('password')
        password_hash = hashlib.sha256(post.password.encode()).hexdigest()

        if stored_password != password_hash:
            raise HTTPException(status_code=403, detail="Invalid password")

        post_ref.update({
            'title': post.title,
            'content': post.content,
            'updatedAt': firestore.SERVER_TIMESTAMP
        })

        return {
            "success": True,
            "message": "Post updated successfully"
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error updating post: {str(e)}")

@app.post("/api/board/posts/{post_id}/verify")
async def verify_post_password(post_id: str, verify: PasswordVerify):
    """
    게시글 비밀번호 검증
    """
    try:
        import hashlib
        db_client = get_firestore_client()
        post_ref = db_client.collection('board_posts').document(post_id)
        post_doc = post_ref.get()

        if not post_doc.exists:
            raise HTTPException(status_code=404, detail="Post not found")

        stored_password = post_doc.to_dict().get('password')
        password_hash = hashlib.sha256(verify.password.encode()).hexdigest()

        return {
            "success": True,
            "valid": stored_password == password_hash
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error verifying password: {str(e)}")

@app.delete("/api/board/posts/{post_id}")
async def delete_board_post(post_id: str, verify: PasswordVerify):
    """
    게시글 삭제
    """
    try:
        import hashlib
        db_client = get_firestore_client()
        post_ref = db_client.collection('board_posts').document(post_id)
        post_doc = post_ref.get()

        if not post_doc.exists:
            raise HTTPException(status_code=404, detail="Post not found")

        # 비밀번호 검증
        stored_password = post_doc.to_dict().get('password')
        password_hash = hashlib.sha256(verify.password.encode()).hexdigest()

        if stored_password != password_hash:
            raise HTTPException(status_code=403, detail="Invalid password")

        post_ref.delete()

        return {
            "success": True,
            "message": "Post deleted successfully"
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting post: {str(e)}")

@app.post("/api/board/posts/{post_id}/like")
async def like_board_post(post_id: str):
    """
    게시글 좋아요
    """
    try:
        db_client = get_firestore_client()
        post_ref = db_client.collection('board_posts').document(post_id)

        if not post_ref.get().exists:
            raise HTTPException(status_code=404, detail="Post not found")

        post_ref.update({
            'likes': firestore.Increment(1)
        })

        # 업데이트된 좋아요 수 가져오기
        updated_post = post_ref.get().to_dict()

        return {
            "success": True,
            "likes": updated_post.get('likes', 1)
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error liking post: {str(e)}")

# ===== 댓글 API =====
@app.get("/api/board/posts/{post_id}/comments")
async def get_comments(post_id: str):
    """
    댓글 목록 조회
    """
    try:
        db_client = get_firestore_client()
        comments_ref = db_client.collection('board_posts').document(post_id).collection('comments')

        query = comments_ref.order_by('createdAt', direction=firestore.Query.ASCENDING)

        comments = []
        for doc in query.stream():
            comment_data = doc.to_dict()
            comment_data['id'] = doc.id
            # 비밀번호는 클라이언트에 전송하지 않음
            comment_data.pop('password', None)
            comments.append(comment_data)

        return {
            "success": True,
            "comments": comments
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching comments: {str(e)}")

@app.post("/api/board/posts/{post_id}/comments")
async def create_comment(post_id: str, comment: BoardComment):
    """
    댓글 작성
    """
    try:
        import hashlib
        db_client = get_firestore_client()

        # 게시글 존재 확인
        post_ref = db_client.collection('board_posts').document(post_id)
        if not post_ref.get().exists:
            raise HTTPException(status_code=404, detail="Post not found")

        comments_ref = post_ref.collection('comments')
        password_hash = hashlib.sha256(comment.password.encode()).hexdigest()

        comment_data = {
            'content': comment.content,
            'author': comment.author,
            'password': password_hash,
            'createdAt': firestore.SERVER_TIMESTAMP,
            'updatedAt': firestore.SERVER_TIMESTAMP
        }

        doc_ref = comments_ref.add(comment_data)
        comment_id = doc_ref[1].id

        return {
            "success": True,
            "message": "Comment created successfully",
            "commentId": comment_id
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating comment: {str(e)}")

@app.put("/api/board/posts/{post_id}/comments/{comment_id}")
async def update_comment(post_id: str, comment_id: str, comment: BoardCommentUpdate):
    """
    댓글 수정
    """
    try:
        import hashlib
        db_client = get_firestore_client()
        comment_ref = db_client.collection('board_posts').document(post_id).collection('comments').document(comment_id)
        comment_doc = comment_ref.get()

        if not comment_doc.exists:
            raise HTTPException(status_code=404, detail="Comment not found")

        # 비밀번호 검증
        stored_password = comment_doc.to_dict().get('password')
        password_hash = hashlib.sha256(comment.password.encode()).hexdigest()

        if stored_password != password_hash:
            raise HTTPException(status_code=403, detail="Invalid password")

        comment_ref.update({
            'content': comment.content,
            'updatedAt': firestore.SERVER_TIMESTAMP
        })

        return {
            "success": True,
            "message": "Comment updated successfully"
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error updating comment: {str(e)}")

@app.delete("/api/board/posts/{post_id}/comments/{comment_id}")
async def delete_comment(post_id: str, comment_id: str, verify: PasswordVerify):
    """
    댓글 삭제
    """
    try:
        import hashlib
        db_client = get_firestore_client()
        comment_ref = db_client.collection('board_posts').document(post_id).collection('comments').document(comment_id)
        comment_doc = comment_ref.get()

        if not comment_doc.exists:
            raise HTTPException(status_code=404, detail="Comment not found")

        # 비밀번호 검증
        stored_password = comment_doc.to_dict().get('password')
        password_hash = hashlib.sha256(verify.password.encode()).hexdigest()

        if stored_password != password_hash:
            raise HTTPException(status_code=403, detail="Invalid password")

        comment_ref.delete()

        return {
            "success": True,
            "message": "Comment deleted successfully"
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting comment: {str(e)}")

# 정적 파일 서빙 (web_service 폴더)
app.mount("/", StaticFiles(directory="web_service", html=True), name="static")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
